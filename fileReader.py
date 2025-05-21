import fitz
from PIL import Image
import pytesseract
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    Docx2txtLoader,
    UnstructuredExcelLoader,
)
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from docx import Document
import os
from os.path import join
from gotOCR import extract_image_data_got_ocr2

pytesseract.pytesseract.tesseract_cmd = (
    "C:\\Program Files\\Tesseract-OCR\\tesseract.exe"
)


def temporary_folder_path():
    folder_path = os.path.join("pipeline", "temporary")
    if not os.path.isdir(folder_path):
        os.makedirs(folder_path)
    return folder_path


def pdf_parser(file_path, file_name):
    text_file_path_list = []
    pdf_file = fitz.open(file_path)

    loader = PyPDFLoader(
        file_path=file_path,
    )

    pdf_txt = loader.load()

    for page_index in range(len(pdf_file)):
        text_per_page = []

        text_per_page.append(pdf_txt[page_index].page_content)

        pdf_page = pdf_file.load_page(page_index)

        images_per_page = pdf_page.get_images(full=True)

        if images_per_page:
            for image_index, img in enumerate(images_per_page, start=1):

                xref = img[0]

                base_image = pdf_file.extract_image(xref)
                image_bytes = base_image["image"]

                image_ext = base_image["ext"]

                image_name = (
                    f"{file_name}_page_{page_index+1}_{image_index}.{image_ext}"
                )
                with open(
                    join(temporary_folder_path(), image_name), "wb"
                ) as image_file:
                    image_file.write(image_bytes)
                    print(f"[+] Image saved as {image_name}")

                # image_text = pytesseract.image_to_string(
                #     join(temporary_folder_path(), image_name)
                # )

                image_text = extract_image_data_got_ocr2(
                    join(temporary_folder_path(), image_name)
                )

                if len(image_text) > 5:
                    text_per_page.append(image_text)

        all_text = "\n".join(text_per_page)

        if len(all_text) > 5:
            text_file_name = f"{file_name}_page_{page_index+1}.txt"
            text_file_path = join(temporary_folder_path(), text_file_name)
            text_file_path_list.append(text_file_path)
            with open(text_file_path, "w") as text_file:
                text_file.write(all_text)

        return text_file_path_list


def pptx_parser(file_path, file_name):
    text_file_path_list = []
    prs = Presentation(file_path)
    image_count = 0

    for idx, slide in enumerate(prs.slides):
        text_per_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_per_slide.append(shape.text)
            elif shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    table_data.append([cell.text for cell in row.cells])
                table_data.append(str(table_data))
            elif shape.shape_type == 13:
                image = shape.image
                ext = image.ext
                image_bytes = image.blob
                image_name = f"{file_name}_slide_{idx+1}_image_{image_count}.{ext}"
                with open(join(temporary_folder_path(), image_name), "wb") as img_file:
                    img_file.write(image_bytes)
                    img_file.close()
                image_count += 1

                image_text = pytesseract.image_to_string(
                    join(temporary_folder_path(), image_name)
                )

                # image_text = extract_image_data_got_ocr2(
                #     join(temporary_folder_path(), image_name)
                # )

                if len(image_text) > 5:
                    text_per_slide.append(image_text)

        all_text = "\n".join(text_per_slide)

        if len(all_text) > 5:
            text_file_name = f"{file_name}_page_{idx+1}.txt"
            text_file_path = join(temporary_folder_path(), text_file_name)
            text_file_path_list.append(text_file_path)
            with open(text_file_path, "w") as text_file:
                text_file.write(all_text)
                text_file.close()

    return text_file_path_list


def create_documents(text_file_path_list):
    document_list = []
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=400,
        chunk_overlap=80,
        length_function=len,
        is_separator_regex=False,
    )
    for i in text_file_path_list:
        with open(i) as f:
            raw_text = f.read()
            document = text_splitter.create_documents([raw_text])
            document_list.extend(document)
            f.close()

    return document_list


def docx_parser(file_path, file_name):
    text_file_path_list = []
    text_per_doc = []
    docx_loader = Docx2txtLoader(file_path)
    text_data = docx_loader.load()

    text_per_doc.append(text_data[0].page_content)

    doc = Document(file_path)
    image_count = 0

    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            image_part = doc.part.rels[rel].target_part
            image_bytes = image_part.blob
            ext = image_part.content_type.split("/")[-1]  # Get file extension
            image_name = f"{file_name}_image_{image_count}.{ext}"
            with open(join(temporary_folder_path(), image_name), "wb") as img_file:
                img_file.write(image_bytes)

            # image_text = pytesseract.image_to_string(
            #     join(temporary_folder_path(), image_name)
            # )

            image_text = extract_image_data_got_ocr2(
                join(temporary_folder_path(), image_name)
            )

            if len(image_text) > 5:
                text_per_doc.append(image_text)

            image_count += 1

    all_text = "\n".join(text_per_doc)

    if len(all_text) > 5:
        text_file_name = f"{file_name}_docx.txt"
        text_file_path = join(temporary_folder_path(), text_file_name)
        text_file_path_list.append(text_file_path)
        with open(text_file_path, "w") as text_file:
            text_file.write(all_text)
            text_file.close()

    return text_file_path_list


def txt_reader(file_path):
    text_loader = TextLoader(file_path)
    text_data = text_loader.load()
    return text_data


def docx_reader(file_path, file_name):
    text_file_path_list = docx_parser(file_path, file_name)
    document_list = create_documents(text_file_path_list)
    return document_list


def pptx_reader(file_path, file_name):
    text_file_path_list = pptx_parser(file_path, file_name)
    document_list = create_documents(text_file_path_list)
    return document_list


def pdf_reader(file_path, file_name):
    text_file_path_list = pdf_parser(file_path, file_name)
    document_list = create_documents(text_file_path_list)
    return document_list


def xlsx_reader(file_path, file_name):
    xlsx_loader = UnstructuredExcelLoader(file_path, mode="elements")
    xlsx_text = xlsx_loader.load()
    return xlsx_text


### Testing ####
# docx_reader(file_path, file_name)
# file_path = "C:\\Users\\skhan\\Documents\\GITHUB\\LANGCHAIN\\Data_Loader\\test_xlsx.xlsx"
# file_name = "test_xlsx.xlsx"

# pptx_reader(file_path, file_name)
# file_path = "C:\\Users\\skhan\\Documents\\GITHUB\\LANGCHAIN\\Data_Loader\\test_docx.docx"
# file_name = "test_docx.docx"


# pdf_reader(file_path, file_name)
# file_path = "C:\\Users\\skhan\\Documents\\GITHUB\\LANGCHAIN\\Data_Loader\\test_data\\test_ppt.pptx"
# file_name = "test_ppt.pptx"
